"""Document parsing service - converts various formats to markdown."""

import io
import json
import logging
from typing import Dict, Optional, Tuple
import PyPDF2
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
from markdownify import markdownify as md

from app.config import Settings

logger = logging.getLogger("tinychat")


class DocumentParser:
    """Parse various document formats to markdown."""
    
    @staticmethod
    async def parse_document(
        file_content: bytes,
        filename: str,
        content_type: str
    ) -> Dict[str, any]:
        """
        Parse document to markdown format.
        
        Args:
            file_content: Raw file bytes
            filename: Original filename
            content_type: MIME type
            
        Returns:
            {
                "markdown": str,     # Parsed content
                "filename": str,     # Original filename
                "size": int,         # File size in bytes
                "type": str,         # Content type
                "pages": int         # Number of pages (if applicable)
            }
            
        Raises:
            ValueError: If file type not supported or parsing fails
        """
        size_mb = len(file_content) / (1024 * 1024)
        if size_mb > Settings.MAX_DOCUMENT_SIZE_MB:
            raise ValueError(
                f"File too large: {size_mb:.1f}MB exceeds limit of {Settings.MAX_DOCUMENT_SIZE_MB}MB"
            )
        
        if content_type not in Settings.SUPPORTED_DOCUMENT_TYPES:
            raise ValueError(f"Unsupported file type: {content_type}")
        
        try:
            # Route to appropriate parser
            if content_type == "text/plain":
                markdown = DocumentParser._parse_text(file_content)
                pages = markdown.count('\n\n') + 1
                
            elif content_type == "text/markdown":
                markdown = file_content.decode('utf-8')
                pages = markdown.count('\n\n') + 1
                
            elif content_type == "text/csv":
                markdown = DocumentParser._parse_csv(file_content)
                pages = 1
                
            elif content_type == "application/pdf":
                markdown, pages = DocumentParser._parse_pdf(file_content)
                
            elif content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                markdown, pages = DocumentParser._parse_docx(file_content)
                
            elif content_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                markdown = DocumentParser._parse_xlsx(file_content)
                pages = 1
                
            elif content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                markdown, pages = DocumentParser._parse_pptx(file_content)
                
            elif content_type == "application/json":
                markdown = DocumentParser._parse_json(file_content)
                pages = 1
                
            elif content_type == "text/html":
                markdown = DocumentParser._parse_html(file_content)
                pages = 1
            
            else:
                raise ValueError(f"Unsupported content type: {content_type}")
            
            logger.info(f"Parsed {filename}: {len(markdown)} chars, {pages} pages")
            
            return {
                "markdown": markdown,
                "filename": filename,
                "size": len(file_content),
                "type": content_type,
                "pages": pages
            }
            
        except Exception as e:
            logger.error(f"Error parsing {filename}: {e}")
            raise ValueError(f"Failed to parse document: {str(e)}")
    
    @staticmethod
    def _parse_text(content: bytes) -> str:
        """Parse plain text file."""
        return content.decode('utf-8', errors='ignore')
    
    @staticmethod
    def _parse_csv(content: bytes) -> str:
        """Parse CSV to markdown table."""
        import csv
        text = content.decode('utf-8', errors='ignore')
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)
        
        if not rows:
            return ""
        
        # Build markdown table
        md_lines = []
        # Header
        md_lines.append("| " + " | ".join(rows[0]) + " |")
        md_lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
        # Data rows
        for row in rows[1:]:
            md_lines.append("| " + " | ".join(row) + " |")
        
        return "\n".join(md_lines)
    
    @staticmethod
    def _parse_pdf(content: bytes) -> Tuple[str, int]:
        """Parse PDF to markdown."""
        pdf_file = io.BytesIO(content)
        reader = PyPDF2.PdfReader(pdf_file)
        
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text.strip():
                pages.append(f"## Page {i + 1}\n\n{text}")
        
        markdown = "\n\n---\n\n".join(pages)
        return markdown, len(reader.pages)
    
    @staticmethod
    def _parse_docx(content: bytes) -> Tuple[str, int]:
        """Parse Word document to markdown."""
        doc_file = io.BytesIO(content)
        doc = DocxDocument(doc_file)
        
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                # Preserve heading styles
                if para.style.name.startswith('Heading'):
                    level = para.style.name.split()[-1]
                    if level.isdigit():
                        paragraphs.append(f"{'#' * int(level)} {para.text}")
                    else:
                        paragraphs.append(f"## {para.text}")
                else:
                    paragraphs.append(para.text)
        
        markdown = "\n\n".join(paragraphs)
        # Estimate pages (rough)
        pages = max(1, len(markdown) // 3000)
        return markdown, pages
    
    @staticmethod
    def _parse_xlsx(content: bytes) -> str:
        """Parse Excel to markdown tables."""
        xlsx_file = io.BytesIO(content)
        wb = load_workbook(xlsx_file, data_only=True)
        
        sheets = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheets.append(f"## Sheet: {sheet_name}\n")
            
            # Convert to markdown table
            rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows.append("| " + " | ".join(str(cell or "") for cell in row) + " |")
            
            if rows:
                # Add header separator after first row
                if len(rows) > 1:
                    header_sep = "| " + " | ".join(["---"] * len(row)) + " |"
                    rows.insert(1, header_sep)
                sheets.append("\n".join(rows))
        
        return "\n\n".join(sheets)
    
    @staticmethod
    def _parse_pptx(content: bytes) -> Tuple[str, int]:
        """Parse PowerPoint to markdown."""
        pptx_file = io.BytesIO(content)
        prs = Presentation(pptx_file)
        
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = [f"## Slide {i + 1}"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if len(slide_text) > 1:  # More than just the header
                slides.append("\n\n".join(slide_text))
        
        markdown = "\n\n---\n\n".join(slides)
        return markdown, len(prs.slides)
    
    @staticmethod
    def _parse_json(content: bytes) -> str:
        """Parse JSON (keep as formatted code block)."""
        json_str = content.decode('utf-8', errors='ignore')
        # Pretty print JSON
        try:
            parsed = json.loads(json_str)
            formatted = json.dumps(parsed, indent=2)
            return f"```json\n{formatted}\n```"
        except:
            # If invalid JSON, return as-is in code block
            return f"```json\n{json_str}\n```"
    
    @staticmethod
    def _parse_html(content: bytes) -> str:
        """Parse HTML to markdown (remove scripts)."""
        html = content.decode('utf-8', errors='ignore')
        soup = BeautifulSoup(html, 'html.parser')
        
        # Remove script and style tags
        for tag in soup(['script', 'style']):
            tag.decompose()
        
        # Convert to markdown
        markdown = md(str(soup), heading_style="ATX")
        return markdown
